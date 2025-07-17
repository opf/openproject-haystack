"""Document processing service for RAG system."""

import os
import logging
from typing import List, Dict, Any, Optional
from pathlib import Path
import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation
import json
from datetime import datetime

logger = logging.getLogger(__name__)


class DocumentChunk:
    """Represents a chunk of text from a document."""
    
    def __init__(
        self,
        text: str,
        metadata: Dict[str, Any],
        chunk_id: str,
        source_file: str,
        page_number: Optional[int] = None
    ):
        self.text = text
        self.metadata = metadata
        self.chunk_id = chunk_id
        self.source_file = source_file
        self.page_number = page_number
        self.created_at = datetime.now().isoformat()


class DocumentProcessor:
    """Processes various document formats for RAG ingestion."""
    
    def __init__(self, chunk_size: int = 2000, chunk_overlap: int = 200):
        """Initialize document processor.
        
        Args:
            chunk_size: Maximum size of text chunks
            chunk_overlap: Number of characters to overlap between chunks
        """
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap
        self.supported_formats = {'.pdf', '.docx', '.pptx', '.txt'}
    
    def process_document(self, file_path: str) -> List[DocumentChunk]:
        """Process a document and return text chunks.
        
        Args:
            file_path: Path to the document file
            
        Returns:
            List of DocumentChunk objects
            
        Raises:
            ValueError: If file format is not supported
            FileNotFoundError: If file doesn't exist
        """
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"Document not found: {file_path}")
        
        file_ext = Path(file_path).suffix.lower()
        if file_ext not in self.supported_formats:
            raise ValueError(f"Unsupported file format: {file_ext}")
        
        logger.info(f"Processing document: {file_path}")
        
        # Extract text based on file type
        if file_ext == '.pdf':
            text_content = self._extract_pdf_text(file_path)
        elif file_ext == '.docx':
            text_content = self._extract_docx_text(file_path)
        elif file_ext == '.pptx':
            text_content = self._extract_pptx_text(file_path)
        elif file_ext == '.txt':
            text_content = self._extract_txt_text(file_path)
        else:
            raise ValueError(f"Unsupported file format: {file_ext}")
        
        logger.info(f"Processing document: {file_path} - before chunks")
        # Create chunks
        chunks = self._create_chunks(text_content, file_path)
        logger.info(f"Processing document: {file_path} - after chunks")
        
        logger.info(f"Created {len(chunks)} chunks from {file_path}")
        return chunks
    
    def _extract_pdf_text(self, file_path: str) -> List[Dict[str, Any]]:
        """Extract text from PDF file.
        
        Args:
            file_path: Path to PDF file
            
        Returns:
            List of dictionaries with text and page information
        """
        text_content = []
        
        try:
            doc = fitz.open(file_path)
            
            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                text = page.get_text()
                
                if text.strip():  # Only add non-empty pages
                    text_content.append({
                        'text': text,
                        'page_number': page_num + 1,
                        'metadata': {
                            'page_count': len(doc),
                            'file_type': 'pdf'
                        }
                    })
            
            doc.close()
            
        except Exception as e:
            logger.error(f"Error extracting PDF text from {file_path}: {e}")
            raise
        
        return text_content
    
    def _extract_docx_text(self, file_path: str) -> List[Dict[str, Any]]:
        """Extract text from Word document.
        
        Args:
            file_path: Path to DOCX file
            
        Returns:
            List of dictionaries with text information
        """
        text_content = []
        
        try:
            doc = Document(file_path)
            
            # Extract paragraphs
            paragraphs = []
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    paragraphs.append(paragraph.text)
            
            # Combine paragraphs into sections
            full_text = '\n\n'.join(paragraphs)
            
            if full_text.strip():
                text_content.append({
                    'text': full_text,
                    'page_number': None,
                    'metadata': {
                        'paragraph_count': len(paragraphs),
                        'file_type': 'docx'
                    }
                })
            
        except Exception as e:
            logger.error(f"Error extracting DOCX text from {file_path}: {e}")
            raise
        
        return text_content
    
    def _extract_pptx_text(self, file_path: str) -> List[Dict[str, Any]]:
        """Extract text from PowerPoint presentation.
        
        Args:
            file_path: Path to PPTX file
            
        Returns:
            List of dictionaries with text and slide information
        """
        text_content = []
        
        try:
            prs = Presentation(file_path)
            
            for slide_num, slide in enumerate(prs.slides):
                slide_text = []
                
                # Extract text from all shapes in the slide
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)
                
                if slide_text:
                    combined_text = '\n\n'.join(slide_text)
                    text_content.append({
                        'text': combined_text,
                        'page_number': slide_num + 1,
                        'metadata': {
                            'slide_count': len(prs.slides),
                            'file_type': 'pptx',
                            'slide_number': slide_num + 1
                        }
                    })
            
        except Exception as e:
            logger.error(f"Error extracting PPTX text from {file_path}: {e}")
            raise
        
        return text_content
    
    def _extract_txt_text(self, file_path: str) -> List[Dict[str, Any]]:
        """Extract text from plain text file.
        
        Args:
            file_path: Path to TXT file
            
        Returns:
            List of dictionaries with text information
        """
        text_content = []
        
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                text = file.read()
                
                if text.strip():
                    text_content.append({
                        'text': text,
                        'page_number': None,
                        'metadata': {
                            'file_type': 'txt',
                            'character_count': len(text)
                        }
                    })
            
        except Exception as e:
            logger.error(f"Error extracting TXT text from {file_path}: {e}")
            raise
        
        return text_content
    
    def _create_chunks(self, text_content: List[Dict[str, Any]], file_path: str) -> List[DocumentChunk]:
        """Create text chunks from extracted content.
        
        Args:
            text_content: List of text content with metadata
            file_path: Source file path
            
        Returns:
            List of DocumentChunk objects
        """
        chunks = []
        chunk_counter = 0
        
        for content in text_content:
            text = content['text']

            page_number = content.get('page_number')
            metadata = content.get('metadata', {})

            
            # Add file-level metadata
            metadata.update({
                'source_file': os.path.basename(file_path),
                'full_path': file_path,
                'processed_at': datetime.now().isoformat()
            })
            
            # Split text into chunks
            text_chunks = self._split_text(text)
            
            for chunk_text in text_chunks:
                chunk_id = f"{os.path.basename(file_path)}_{chunk_counter}"
                
                chunk = DocumentChunk(
                    text=chunk_text,
                    metadata=metadata.copy(),
                    chunk_id=chunk_id,
                    source_file=file_path,
                    page_number=page_number
                )
                
                chunks.append(chunk)
                chunk_counter += 1
        
        return chunks
    
    def _split_text(self, text: str) -> List[str]:
        """Split text into chunks with overlap.
        
        Args:
            text: Text to split
            
        Returns:
            List of text chunks
        """
        if len(text) <= self.chunk_size:
            return [text]
        
        chunks = []
        start = 0
        iteration_count = 0
        max_iterations = len(text) // 10  # Safety limit to prevent infinite loops
        
        while start < len(text) and iteration_count < max_iterations:
            iteration_count += 1
            original_start = start
            end = min(start + self.chunk_size, len(text))
            
            logger.debug(f"Iteration {iteration_count}: start={start}, end={end}, text_len={len(text)}")
            
            # If this is not the last chunk, try to break at a good boundary
            if end < len(text):
                # Look for sentence boundary (. ! ?) within a reasonable range
                search_start = max(start, end - 200)  # Don't search too far back
                sentence_end = max(
                    text.rfind('.', search_start, end),
                    text.rfind('!', search_start, end),
                    text.rfind('?', search_start, end)
                )
                
                if sentence_end > start:
                    end = sentence_end + 1
                else:
                    # Look for word boundary within a reasonable range
                    search_start = max(start, end - 100)  # Don't search too far back
                    word_end = text.rfind(' ', search_start, end)
                    if word_end > start:
                        end = word_end
                    # If no good boundary found, use the original end position
            
            # Extract chunk
            chunk = text[start:end].strip()
            if chunk:
                chunks.append(chunk)
                logger.debug(f"Created chunk {len(chunks)}: length={len(chunk)}")
            
            # Calculate next start position with overlap
            # Ensure we always make forward progress
            next_start = max(end - self.chunk_overlap, start + 1)
            
            # Additional safety check: if we're not making progress, force advancement
            if next_start <= start:
                next_start = start + max(1, self.chunk_size // 2)
                logger.warning(f"Forced advancement from {start} to {next_start} to prevent infinite loop")
            
            start = next_start
            
            # Safety check for infinite loop
            if start == original_start and iteration_count > 1:
                logger.error(f"Infinite loop detected at position {start}. Forcing advancement.")
                start = original_start + max(1, self.chunk_size // 2)
        
        if iteration_count >= max_iterations:
            logger.error(f"Maximum iterations ({max_iterations}) reached. Text may not be fully processed.")
        
        logger.info(f"Split text into {len(chunks)} chunks in {iteration_count} iterations")
        return chunks
    
    def get_document_metadata(self, file_path: str) -> Dict[str, Any]:
        """Get metadata for a document without processing full content.
        
        Args:
            file_path: Path to document
            
        Returns:
            Dictionary with document metadata
        """
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"Document not found: {file_path}")
        
        file_stat = os.stat(file_path)
        file_ext = Path(file_path).suffix.lower()
        
        metadata = {
            'filename': os.path.basename(file_path),
            'file_path': file_path,
            'file_extension': file_ext,
            'file_size': file_stat.st_size,
            'modified_time': datetime.fromtimestamp(file_stat.st_mtime).isoformat(),
            'created_time': datetime.fromtimestamp(file_stat.st_ctime).isoformat()
        }
        
        # Add format-specific metadata
        try:
            if file_ext == '.pdf':
                doc = fitz.open(file_path)
                metadata.update({
                    'page_count': len(doc),
                    'title': doc.metadata.get('title', ''),
                    'author': doc.metadata.get('author', ''),
                    'subject': doc.metadata.get('subject', '')
                })
                doc.close()
            
            elif file_ext == '.docx':
                doc = Document(file_path)
                metadata.update({
                    'paragraph_count': len(doc.paragraphs),
                    'title': doc.core_properties.title or '',
                    'author': doc.core_properties.author or '',
                    'subject': doc.core_properties.subject or ''
                })
            
            elif file_ext == '.pptx':
                prs = Presentation(file_path)
                metadata.update({
                    'slide_count': len(prs.slides),
                    'title': prs.core_properties.title or '',
                    'author': prs.core_properties.author or '',
                    'subject': prs.core_properties.subject or ''
                })
        
        except Exception as e:
            logger.warning(f"Could not extract detailed metadata from {file_path}: {e}")
        
        return metadata
